# -*- coding: utf-8 -*-
"""文書の要約機能（AI・Gemini経由）── 要約機能 Phase A/B

依頼:「検索結果の内容を、①300字のExecutive Summary ②章立てと概要
③Japan Site Managerへの示唆、の3段階でAI要約したい」への対応。

対象は SharePoint / Nexus の .docx / .pptx（Phase Aで.docx、Phase Bで.pptx
を追加）。フォルダ・Enovia・他形式（.pdf/.xlsx）は対象外（設計議論で合意済み、
DESIGN_NOTES参照）。AI呼び出しは既存の翻訳ツール群と同じ共通モジュール
gemini_client.py を流用するが、本テストはネットワークには一切アクセスせず、
_generate_advanced 等をスタブに差し替えて検証する。

実行: python tests/test_17_summary.py
"""
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
from _harness import dsm, DummyAuth  # noqa: E402

import io
import json

try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

ok, ng = 0, 0


def check(label, cond, detail=""):
    global ok, ng
    if cond:
        ok += 1
        print(f"  OK   {label}")
    else:
        ng += 1
        print(f"  NG   {label}  {detail}")


def make_docx_bytes(paragraphs):
    """[(text, style_or_None), ...] からdocxのバイト列を作る。"""
    doc = Document()
    for text, style in paragraphs:
        if style:
            doc.add_paragraph(text, style=style)
        else:
            doc.add_paragraph(text)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def make_pptx_bytes(slides):
    """[(title_or_None, [本文の段落, ...]), ...] からpptxのバイト列を作る。"""
    from pptx.util import Inches

    prs = Presentation()
    layout_with_title = prs.slide_layouts[1]   # タイトル＋本文
    layout_blank = prs.slide_layouts[6]        # 白紙（タイトルプレースホルダ無し）

    for title, body_lines in slides:
        if title is not None:
            slide = prs.slides.add_slide(layout_with_title)
            slide.shapes.title.text = title
            if body_lines:
                body_tf = slide.placeholders[1].text_frame
                body_tf.text = body_lines[0]
                for line in body_lines[1:]:
                    body_tf.add_paragraph().text = line
        else:
            slide = prs.slides.add_slide(layout_blank)
            if body_lines:
                box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(3))
                tf = box.text_frame
                tf.text = body_lines[0]
                for line in body_lines[1:]:
                    tf.add_paragraph().text = line

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def fake_gemini_response(executive_summary="要約です。",
                         chapters=None, insights=None):
    chapters = chapters if chapters is not None else [
        {"title": "第1章 概要", "overview": "概要の説明。"},
    ]
    insights = insights if insights is not None else {
        "use": ["活用点1", "活用点2", "活用点3"],
        "caution": ["注意点1", "注意点2", "注意点3"],
        "questions": ["問い1", "問い2", "問い3"],
    }
    data = {"executive_summary": executive_summary, "chapters": chapters,
            "insights": insights}
    return {"candidates": [{"content": {"parts": [
        {"text": json.dumps(data, ensure_ascii=False)}
    ]}}]}


CFG = dict(dsm.DEFAULT_CFG)


# ── G1 _summarizable_reason（Phase Aの対象を絞るガード） ────────
print("\n[G1] 要約可能かどうかの判定（_summarizable_reason）")


def mk(**kwargs):
    base = dict(source="SharePoint", title="Sample", doc_type="docx",
                url="https://x/sites/S/Docs/Sample.docx", last_modified="2026-09-01")
    base.update(kwargs)
    return dsm.SearchResult(**base)


ok_row = mk()
check("SharePointのdocxは要約可能", dsm._summarizable_reason(ok_row) == (True, ""))

folder_row = mk(is_folder=True, doc_type=dsm.FOLDER_TYPE_LABEL)
reason = dsm._summarizable_reason(folder_row)
check("フォルダは要約不可", reason[0] is False and "フォルダ" in reason[1], reason)

enovia_row = mk(source="Enovia", doc_type="docx")
reason = dsm._summarizable_reason(enovia_row)
check("Enoviaは要約不可（Phase A対象外）", reason[0] is False and "Enovia" in reason[1], reason)

no_url_row = mk(url="")
reason = dsm._summarizable_reason(no_url_row)
check("リンクが無い行は要約不可", reason[0] is False and "リンク" in reason[1], reason)

pptx_row = mk(doc_type="pptx")
check("pptxも要約可能（Phase Bで追加）",
      dsm._summarizable_reason(pptx_row) == (True, ""), dsm._summarizable_reason(pptx_row))

pdf_row = mk(doc_type="pdf")
reason = dsm._summarizable_reason(pdf_row)
check("docx/pptx以外（pdf）は現状要約不可", reason[0] is False and "pdf" in reason[1], reason)

nexus_row = mk(source="Nexus")
check("Nexusのdocxも要約可能（SharePointProvider系統を継承）",
      dsm._summarizable_reason(nexus_row) == (True, ""))


# ── G2 _extract_docx_text（本文抽出・見出し・切り詰め） ─────────
if HAS_DOCX:
    print("\n[G2] docx本文抽出（_extract_docx_text）")

    content = make_docx_bytes([
        ("第1章 はじめに", "Heading 1"),
        ("これは本文の1段落目です。", None),
        ("第2章 詳細", "Heading 1"),
        ("これは本文の2段落目です。", None),
        ("", None),   # 空段落は無視される
    ])
    text, truncated = dsm._extract_docx_text(content, max_chars=10000)
    check("見出し段落に # が付く", "# 第1章 はじめに" in text, text)
    check("本文段落はそのまま", "これは本文の1段落目です。" in text, text)
    check("空段落は含まれない（連続する空行にならない）",
          "\n\n\n" not in text, repr(text))
    check("十分に短ければ切り詰められない", truncated is False)

    long_content = make_docx_bytes([("あ" * 100, None)])
    long_text, long_truncated = dsm._extract_docx_text(long_content, max_chars=50)
    check("上限を超えると切り詰められる", long_truncated is True)
    check("切り詰め後の長さが上限以下", len(long_text) <= 50, len(long_text))

    empty_content = make_docx_bytes([("", None)])
    empty_text, _ = dsm._extract_docx_text(empty_content, max_chars=10000)
    check("本文が無ければ空文字", empty_text == "", repr(empty_text))
else:
    print("\n[G2] python-docx未インストールのためスキップ（他機能には影響しない設計）")


# ── G2b _extract_pptx_text（本文抽出・スライド区切り・切り詰め） ────
if HAS_PPTX:
    print("\n[G2b] pptx本文抽出（_extract_pptx_text）── Phase Bで追加")

    content = make_pptx_bytes([
        ("スライド1のタイトル", ["本文1行目", "本文2行目"]),
        (None, ["タイトル無しスライドの本文"]),
    ])
    text, truncated = dsm._extract_pptx_text(content, max_chars=10000)
    check("スライドの区切りにタイトルが付く", "# Slide 1: スライド1のタイトル" in text, text)
    check("スライド本文が含まれる", "本文1行目" in text and "本文2行目" in text, text)
    check("タイトル無しスライドは番号のみの区切りになる", "# Slide 2" in text, text)
    check("タイトル無しスライドの本文も含まれる",
          "タイトル無しスライドの本文" in text, text)
    check("十分に短ければ切り詰められない", truncated is False)

    long_content = make_pptx_bytes([("長いスライド", ["あ" * 100])])
    long_text, long_truncated = dsm._extract_pptx_text(long_content, max_chars=50)
    check("上限を超えると切り詰められる", long_truncated is True)
    check("切り詰め後の長さが上限以下", len(long_text) <= 50, len(long_text))

    print("\n[G2c] 拡張子による振り分け（_extract_text_for_summary）")
    docx_bytes_for_dispatch = make_docx_bytes([("docxの本文", None)])
    pptx_bytes_for_dispatch = make_pptx_bytes([("見出し", ["pptxの本文"])])
    d_text, _ = dsm._extract_text_for_summary("docx", docx_bytes_for_dispatch, 10000)
    p_text, _ = dsm._extract_text_for_summary("pptx", pptx_bytes_for_dispatch, 10000)
    check("doc_type=docxならdocx抽出が使われる", "docxの本文" in d_text, d_text)
    check("doc_type=pptxならpptx抽出が使われる", "pptxの本文" in p_text, p_text)
else:
    print("\n[G2b] python-pptx未インストールのためスキップ（他機能には影響しない設計）")


# ── G3 _generate_summary（Geminiレスポンスの解析） ──────────────
print("\n[G3] Geminiレスポンスの解析（_generate_summary）")

captured = {}


def stub_ok(payload, model=None):
    captured["payload"] = payload
    captured["model"] = model
    return fake_gemini_response()


orig_generate_advanced = dsm._generate_advanced
dsm._generate_advanced = stub_ok
try:
    result = dsm._generate_summary("Sample", "本文テキスト")
    check("executive_summaryを取得できる", result["executive_summary"] == "要約です。", result)
    check("chaptersを取得できる", result["chapters"][0]["title"] == "第1章 概要", result)
    check("insightsを取得できる（箇条書きの配列）",
          result["insights"]["use"] == ["活用点1", "活用点2", "活用点3"], result)
    check("モデル名が渡される", captured["model"] == dsm.GEMINI_MODEL_NAME, captured.get("model"))
    check("responseMimeTypeがJSON指定", captured["payload"]["generationConfig"]["responseMimeType"]
          == "application/json")
    check("responseSchemaが渡される",
          "responseSchema" in captured["payload"]["generationConfig"])

    dsm._generate_advanced = lambda payload, model=None: {"candidates": []}
    try:
        dsm._generate_summary("Sample", "本文")
        check("candidatesが空なら例外", False, "例外が発生しなかった")
    except RuntimeError as e:
        check("candidatesが空なら例外", "応答" in str(e), str(e))

    dsm._generate_advanced = lambda payload, model=None: {
        "candidates": [{"content": {"parts": [{"text": "not-json"}]}}]}
    try:
        dsm._generate_summary("Sample", "本文")
        check("JSONでない応答は例外", False, "例外が発生しなかった")
    except RuntimeError as e:
        check("JSONでない応答は例外", "JSON" in str(e), str(e))

    dsm._generate_advanced = lambda payload, model=None: {"candidates": [{"content": {"parts": [
        {"text": json.dumps({"executive_summary": "x"})}
    ]}}]}
    try:
        dsm._generate_summary("Sample", "本文")
        check("必須キー欠落は例外", False, "例外が発生しなかった")
    except RuntimeError as e:
        check("必須キー欠落は例外", "chapters" in str(e), str(e))
finally:
    dsm._generate_advanced = orig_generate_advanced


# ── G4 /api/summarize エンドポイント ────────────────────────────
print("\n[G4] /api/summarize エンドポイント")


def idx_of(resp_json, title):
    """/api/search の応答から、タイトルで行の_idx相当（配列位置）を探す。

    SearchManager.search() は結果を (source, rank) でソートするため、
    渡した順序どおりに並ぶとは限らない（実際にEnovia行が先頭に来ることを
    本テスト作成時に確認済み）。そのためタイトルで引き直す。
    """
    for i, r in enumerate(resp_json["results"]):
        if r["title"] == title:
            return str(i)
    raise AssertionError(f"title={title} が見つかりません: {resp_json}")


dsm._cfg = CFG
dsm._auth = DummyAuth()
mgr = dsm.SearchManager(CFG, DummyAuth())
rows = [
    mk(title="DocA", url="https://x/sites/S/Docs/A.docx",
       last_modified="2026-09-01"),                                             # 要約可
    mk(title="FolderX", is_folder=True, doc_type=dsm.FOLDER_TYPE_LABEL,
       url="https://x/sites/S/Docs/Folder"),                                    # フォルダ
    mk(title="EnoviaDoc", source="Enovia", url="https://x/sites/S/Docs/B.docx"),  # Enovia
    mk(title="PdfOne", doc_type="pdf", url="https://x/sites/S/Docs/C.pdf"),      # 未対応形式
    mk(title="PptxDoc", url="https://x/sites/S/Docs/D.pptx", doc_type="pptx",
       last_modified="2026-09-01"),                                             # 要約可（Phase B）
]
mgr.providers[dsm.TARGET_SHAREPOINT].search = lambda kw, mx: {
    "results": list(rows), "total": len(rows), "note": ""}
dsm._manager = mgr
client = dsm.flask_app.test_client()
search_resp = client.post("/api/search", json={"keyword": "sample", "target": "sharepoint"}).get_json()
idx_doc_a = idx_of(search_resp, "DocA")
idx_pptx_doc = idx_of(search_resp, "PptxDoc")

r = client.post("/api/summarize", json={})
check("idx未指定は400", r.status_code == 400, r.status_code)

r = client.post("/api/summarize", json={"idx": idx_of(search_resp, "FolderX")})
check("フォルダ行は400（理由付き）",
      r.status_code == 400 and "フォルダ" in r.get_json()["error"], r.get_json())

r = client.post("/api/summarize", json={"idx": idx_of(search_resp, "EnoviaDoc")})
check("Enovia行は400（理由付き）",
      r.status_code == 400 and "Enovia" in r.get_json()["error"], r.get_json())

r = client.post("/api/summarize", json={"idx": idx_of(search_resp, "PdfOne")})
check("pdf行は400（理由付き・現状docx/pptxのみ対応）",
      r.status_code == 400 and "pdf" in r.get_json()["error"], r.get_json())

r = client.post("/api/summarize", json={"idx": "999"})
check("範囲外の索引は400", r.status_code == 400, r.status_code)

# Gemini未導入・未設定の状態からの検証（このサンドボックスの既定状態）
orig_has_gemini = dsm.HAS_GEMINI
orig_cred_check = dsm.gemini_credentials_available
if not HAS_DOCX:
    print("  (python-docx未インストールのため、以降のG4検証はスキップ)")
else:
    dsm.HAS_GEMINI = False
    r = client.post("/api/summarize", json={"idx": idx_doc_a})
    check("Gemini共通モジュール未導入なら500", r.status_code == 500, r.status_code)

    dsm.HAS_GEMINI = True
    dsm.gemini_credentials_available = lambda: False
    r = client.post("/api/summarize", json={"idx": idx_doc_a})
    check("認証情報未設定なら500", r.status_code == 500
          and "GEMINI_API_KEY" in r.get_json()["error"], r.get_json())
    dsm.gemini_credentials_available = lambda: True

    # ダウンロード・Gemini呼び出しをスタブ化した正常系
    docx_bytes = make_docx_bytes([
        ("概要", "Heading 1"),
        ("これはテスト用の十分に長い本文です。" * 5, None),
    ])

    class FakeResp:
        def __init__(self, code, content=b""):
            self.status_code = code
            self.content = content

    download_calls = []

    def fake_get(url, headers=None, timeout=None, allow_redirects=None):
        download_calls.append(url)
        return FakeResp(200, docx_bytes)

    orig_http_get = dsm.http_req.get
    dsm.http_req.get = fake_get

    gemini_calls = []

    def counting_gemini(payload, model=None):
        gemini_calls.append(payload)
        return fake_gemini_response(executive_summary="これはExecutive Summaryです。")

    dsm._generate_advanced = counting_gemini

    try:
        r = client.post("/api/summarize", json={"idx": idx_doc_a})
        check("正常系は200", r.status_code == 200, r.status_code)
        data = r.get_json()
        check("executive_summaryが返る",
              data.get("executive_summary") == "これはExecutive Summaryです。", data)
        check("chaptersが返る", len(data.get("chapters", [])) == 1, data)
        check("insightsが返る", "use" in data.get("insights", {}), data)
        check("titleが返る", data.get("title") == "DocA", data)
        check("truncatedはFalse（十分短いため）", data.get("truncated") is False, data)
        check("ダウンロードは1回呼ばれる", len(download_calls) == 1, len(download_calls))
        check("Geminiは1回呼ばれる", len(gemini_calls) == 1, len(gemini_calls))

        # ── キャッシュ: 同じ条件なら再取得・再要約しない ──
        r2 = client.post("/api/summarize", json={"idx": idx_doc_a})
        check("2回目もキャッシュにより200", r2.status_code == 200, r2.status_code)
        check("キャッシュにより再ダウンロードしない",
              len(download_calls) == 1, len(download_calls))
        check("キャッシュにより再度Geminiを呼ばない",
              len(gemini_calls) == 1, len(gemini_calls))

        # ── 文書が更新されたら（last_modifiedが変われば）キャッシュを使わない ──
        rows[0].last_modified = "2026-09-02"
        mgr.providers[dsm.TARGET_SHAREPOINT].search = lambda kw, mx: {
            "results": list(rows), "total": len(rows), "note": ""}
        client.post("/api/search", json={"keyword": "sample", "target": "sharepoint"})
        r3 = client.post("/api/summarize", json={"idx": idx_doc_a})
        check("最終更新日が変わればキャッシュを使わず再取得する",
              len(download_calls) == 2, len(download_calls))
        check("最終更新日が変わればキャッシュを使わず再度要約する",
              len(gemini_calls) == 2, len(gemini_calls))

        # ── 本文がほとんど無い場合は要約しない（誤った要約を避けるガード） ──
        dsm.http_req.get = lambda *a, **k: FakeResp(200, make_docx_bytes([("短", None)]))
        rows[0].last_modified = "2026-09-03"
        mgr.providers[dsm.TARGET_SHAREPOINT].search = lambda kw, mx: {
            "results": list(rows), "total": len(rows), "note": ""}
        client.post("/api/search", json={"keyword": "sample", "target": "sharepoint"})
        r4 = client.post("/api/summarize", json={"idx": idx_doc_a})
        check("本文が短すぎる場合は422（要約を行わない）",
              r4.status_code == 422, (r4.status_code, r4.get_json()))
        check("422のときGeminiは呼ばれない（無駄打ちしない）",
              len(gemini_calls) == 2, len(gemini_calls))

        # ── 切り詰め（summary_max_charsで上限を絞る） ──
        dsm.http_req.get = fake_get
        dsm._cfg = dict(CFG, summary_max_chars=20, summary_min_chars=5)
        rows[0].last_modified = "2026-09-04"
        mgr.providers[dsm.TARGET_SHAREPOINT].search = lambda kw, mx: {
            "results": list(rows), "total": len(rows), "note": ""}
        client.post("/api/search", json={"keyword": "sample", "target": "sharepoint"})
        r5 = client.post("/api/summarize", json={"idx": idx_doc_a})
        check("上限文字数を絞るとtruncated=True", r5.get_json().get("truncated") is True,
              r5.get_json())
        dsm._cfg = CFG

        # ── ダウンロード失敗時のエラー ──
        rows[0].last_modified = "2026-09-05"
        mgr.providers[dsm.TARGET_SHAREPOINT].search = lambda kw, mx: {
            "results": list(rows), "total": len(rows), "note": ""}
        client.post("/api/search", json={"keyword": "sample", "target": "sharepoint"})
        dsm.http_req.get = lambda *a, **k: FakeResp(403)
        r6 = client.post("/api/summarize", json={"idx": idx_doc_a})
        check("ファイル取得失敗は500", r6.status_code == 500, r6.status_code)

        # ── Gemini呼び出し失敗時のエラー ──
        rows[0].last_modified = "2026-09-06"
        mgr.providers[dsm.TARGET_SHAREPOINT].search = lambda kw, mx: {
            "results": list(rows), "total": len(rows), "note": ""}
        client.post("/api/search", json={"keyword": "sample", "target": "sharepoint"})
        dsm.http_req.get = fake_get

        def failing_gemini(payload, model=None):
            raise RuntimeError("プロキシに到達できません")

        dsm._generate_advanced = failing_gemini
        r7 = client.post("/api/summarize", json={"idx": idx_doc_a})
        check("Gemini呼び出し失敗は502", r7.status_code == 502
              and "要約の生成に失敗" in r7.get_json()["error"], r7.get_json())

        # ── pptx（Phase Bで追加）でも /api/summarize が一連で動く ──
        if HAS_PPTX:
            pptx_bytes = make_pptx_bytes([
                ("スライド1", ["これはpptxのテスト用の本文です。" * 5]),
            ])
            dsm.http_req.get = lambda *a, **k: FakeResp(200, pptx_bytes)
            dsm._generate_advanced = counting_gemini
            gemini_calls_before = len(gemini_calls)
            r8 = client.post("/api/summarize", json={"idx": idx_pptx_doc})
            check("pptxでも正常系は200", r8.status_code == 200, r8.status_code)
            check("pptxでもGeminiが呼ばれる（拡張子で抽出関数が振り分けられている）",
                  len(gemini_calls) == gemini_calls_before + 1, len(gemini_calls))
        else:
            print("  (python-pptx未インストールのため、pptxのエンドツーエンド検証はスキップ)")
    finally:
        dsm.http_req.get = orig_http_get
        dsm._generate_advanced = orig_generate_advanced
        dsm.HAS_GEMINI = orig_has_gemini
        dsm.gemini_credentials_available = orig_cred_check
        dsm._cfg = CFG


# ── G5 画面表示（ボタン・ポップアップの存在確認） ────────────────
print("\n[G5] 画面表示")
html = dsm.flask_app.test_client().get("/").get_data(as_text=True)
check("要約列がAllタブに定義されている", '"__summary"' in html)
check("要約ボタンの生成ロジックがある", "openSummaryPopup" in html)
check("要約可否の判定ロジックがある（Phase Aの対象を絞る）", "summaryUnavailableReason" in html)
check("ポップアップのDOM構造がある", "summaryOverlay" in html and "summary-modal" in html)
check("/api/summarize を呼び出すfetchがある", '"/api/summarize"' in html)
check("Escキーで閉じるハンドラがある", "summaryPopupKeyHandler" in html)
check("示唆を箇条書き（配列）で描画するロジックがある（文章と箇条書きの混在対策）",
      "summary-insight-list" in html)
check("画面側の対応形式判定にpptxが含まれる（Phase B）",
      'SUMMARIZABLE_EXTENSIONS = ["docx", "pptx"]' in html)


print(f"\n{'=' * 46}\n  成功 {ok} 件 / 失敗 {ng} 件\n{'=' * 46}")
sys.exit(1 if ng else 0)
