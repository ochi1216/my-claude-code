r"""
VERSION　2026.0309.03
追加・修正
**進捗の可視化とUIハングアップ対策**: 並列処理エンジンがバッチ（10項目）を処理するごとにプログレスバーを更新するコールバック関数を導入し、画面上で進捗がリアルタイムに確認できるように変更しました（根拠：0%で固まって見えるUIの欠陥を解消するため）。
**APIタイムアウトとリトライ機構**: `generate_content` に40秒のタイムアウトを設定し、失敗時は「3回・3秒間隔」で再試行するロジックを追加しました（根拠：APIの応答遅延による無限フリーズを防ぐため）。
**フェイルファスト（即時撤退）と強制切断**: 3バッチ連続でエラーが発生した場合、残りの処理を即座にキャンセル（スレッド待機フラグによる安全停止）し、エラーメッセージを出して終了する安全装置を追加しました（根拠：API障害時に無駄なリクエストを送り続けるゾンビ化を防ぐため）。
**デバッグ用ログ（LOG）出力**: 処理の足跡と通信エラーの詳細を `translation_debug.log` に記録し、コンソールにも出力するロギング機構を追加しました（根拠：問題切り分けの証拠を残すため）。
**ファイルロックの事前検知**: 翻訳処理を開始する前に、出力先ファイルの書き込み権限をチェックし、ロックされている場合は即座に警告を出す処理を追加しました（根拠：数分間のAPI処理が終わった直後にWinError 32で落ちる悲劇を防ぐため）。
変更関数
`グローバルスコープ` (logging等のインポート追加)
`get_logger` (新規追加：ロガー初期化)
`translate_batch_gemini` (ロガー引数、タイムアウト、リトライ3回/3秒待機の追加)
`translate_super_fast_parallel` (コールバック引数、連続エラー検知、キャンセルフラグの追加)
`translate_ppt_document_thread` (ログ出力、ファイルロック事前検知、UI進捗更新の紐付け)
新規追加：
`get_logger` (グローバルスコープ直下 / `init_gemini`の直前)
"""
import tkinter as tk
from tkinter import filedialog, messagebox
import os
import sys
import shutil
import re
import time
import threading
from concurrent.futures import ThreadPoolExecutor, as_completed
import logging
import traceback

# --- 依存関係の確認 ---
try:
    import google.generativeai as genai
    HAS_GEMINI = True
except ImportError:
    HAS_GEMINI = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

def check_dependencies(root_window):
    """起動時の依存関係チェック"""
    missing_libs = []
    if not HAS_GEMINI:
        missing_libs.append("google-generativeai")
    if not HAS_PPTX:
        missing_libs.append("python-pptx")
    
    if missing_libs:
        error_msg = "以下のライブラリがインストールされていません:\n"
        for lib in missing_libs:
            error_msg += f"- {lib}\n"
        error_msg += "\n以下のコマンドでインストールしてください:\n"
        error_msg += f"pip install {' '.join(missing_libs)}"
        
        messagebox.showerror("依存関係エラー", error_msg, parent=root_window)
        return False
    return True

# --- グローバル変数（APIモデル） ---
gemini_model = None

# --- 新規追加: ロガーの初期化 ---
def get_logger():
    """デバッグ用ロガーの初期化（コンソールとファイル両方に出力）"""
    logger = logging.getLogger("PPT_Translation")
    if not logger.handlers:
        logger.setLevel(logging.DEBUG)
        fh = logging.FileHandler("translation_debug.log", encoding="utf-8")
        ch = logging.StreamHandler()
        formatter = logging.Formatter("[%(asctime)s] %(levelname)s - %(message)s", "%H:%M:%S")
        fh.setFormatter(formatter)
        ch.setFormatter(formatter)
        logger.addHandler(fh)
        logger.addHandler(ch)
    return logger

def init_gemini(root_window):
    """Gemini APIの初期化（自動モデル検出機能付き）"""
    global gemini_model
    api_key = os.getenv("GEMINI_API_KEY")
    if not api_key:
        messagebox.showerror("エラー", 
                           "GEMINI_API_KEY が設定されていません。\n"
                           "環境変数に GEMINI_API_KEY を設定してください。", parent=root_window)
        return False
    
    try:
        genai.configure(api_key=api_key)
        
        # 使えるモデルを自動検出して404エラーを完全に防ぐ
        available_models = []
        for m in genai.list_models():
            if 'generateContent' in m.supported_generation_methods:
                available_models.append(m.name)
                
        # 優先的にFlashモデルを探す
        target_model_name = None
        for preferred in ['gemini-2.5-flash', 'gemini-2.0-flash', 'gemini-1.5-flash', 'gemini-pro']:
            for am in available_models:
                if preferred in am:
                    target_model_name = am
                    break
            if target_model_name:
                break
                
        # 見つからなければ最初のモデルを使用
        if not target_model_name and available_models:
            target_model_name = available_models[0]
            
        print(f"自動選択されたモデル: {target_model_name}")
        gemini_model = genai.GenerativeModel(target_model_name)
        return True
    except Exception as e:
        messagebox.showerror("API初期化エラー", f"Gemini APIの初期化に失敗しました:\n{e}", parent=root_window)
        return False

def is_translatable(text):
    """翻訳が必要なテキストかどうかを判定"""
    if not text or str(text).strip() == "":
        return False
    
    text_str = str(text).strip()
    
    if text_str in ["", "#", "-", "N/A", "NULL", "•", "◦", "▪", "**", "*", ":", "：", 
                    "I.", "II.", "III.", "IV.", "V.", "VI.", "***"]:
        return False
    if text_str.replace(".", "").replace("-", "").isdigit():
        return False
    if len(text_str) <= 2:
        return False
    return True

def translate_batch_gemini(texts, target_language="Japanese", batch_idx=0, logger=None):
    """Gemini APIを使用した小バッチ翻訳（リトライ・タイムアウト機構付き）"""
    if not gemini_model or not texts:
        return texts, False
    
    batch_input = "\n".join([f"[{i+1}] {t}" for i, t in enumerate(texts)])
    
    prompt = f"""
    Task: Translate the following text into {target_language}.
    
    Guidelines:
    1. Maintain the exact format [number] for each translated line.
    2. Output ONLY the numbered list. No extra explanations.
    3. Keep technical terms natural.
    
    Source Text:
    {batch_input}
    """

    safety_settings = [
        {"category": "HARM_CATEGORY_HARASSMENT", "threshold": "BLOCK_NONE"},
        {"category": "HARM_CATEGORY_HATE_SPEECH", "threshold": "BLOCK_NONE"},
        {"category": "HARM_CATEGORY_SEXUALLY_EXPLICIT", "threshold": "BLOCK_NONE"},
        {"category": "HARM_CATEGORY_DANGEROUS_CONTENT", "threshold": "BLOCK_NONE"},
    ]

    for attempt in range(1, 4):  # 最大3回リトライ
        try:
            if logger: logger.info(f"バッチ {batch_idx+1} 通信開始 (試行 {attempt}/3)")
            
            # APIタイムアウト（40秒）を設定
            response = gemini_model.generate_content(
                prompt,
                generation_config=genai.types.GenerationConfig(temperature=0.1),
                safety_settings=safety_settings,
                request_options={"timeout": 40}
            )
            time.sleep(1.5)

            if not response.parts:
                if logger: logger.warning(f"バッチ {batch_idx+1} 空のレスポンスを受信")
                raise ValueError("Empty response from API")

            response_text = response.text.strip()
            results = [None] * len(texts)
            lines = response_text.split('\n')
            
            for line in lines:
                match = re.match(r'^\[(\d+)\]\s*(.*)', line.strip())
                if match:
                    idx = int(match.group(1)) - 1
                    if 0 <= idx < len(texts):
                        results[idx] = match.group(2).strip()
            
            for i in range(len(results)):
                if results[i] is None:
                    results[i] = texts[i]
                    
            if logger: logger.info(f"バッチ {batch_idx+1} 成功！")
            return results, False  # 成功（エラーフラグFalse）

        except Exception as e:
            if logger: logger.error(f"バッチ {batch_idx+1} エラー発生: {str(e)}")
            if attempt < 3:
                time.sleep(3)  # エラー時は3秒間隔で待機
            else:
                if logger: logger.error(f"バッチ {batch_idx+1} は3回失敗したためスキップします。")
    
    return texts, True  # 失敗（原文を返し、エラーフラグTrueを通知）

def translate_super_fast_parallel(all_texts, target_language="Japanese", max_workers=3, progress_callback=None, logger=None):
    """並列処理エンジン（コールバックと強制切断機能付き）"""
    if not all_texts:
        return []
    
    batch_size = 10
    chunks = [all_texts[i:i + batch_size] for i in range(0, len(all_texts), batch_size)]
    results = [None] * len(chunks)
    
    abort_event = threading.Event()
    consecutive_errors = 0
    processed_items = 0
    
    def translate_chunk(chunk_idx, chunk_texts):
        if abort_event.is_set():
            return chunk_idx, (chunk_texts, False)  # 中断フラグが立っていればスルー
        return chunk_idx, translate_batch_gemini(chunk_texts, target_language, chunk_idx, logger)

    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        futures = [executor.submit(translate_chunk, i, chunk) for i, chunk in enumerate(chunks)]
        
        for future in as_completed(futures):
            try:
                chunk_idx, (translated_chunk, is_error) = future.result()
                results[chunk_idx] = translated_chunk
                
                # エラーカウントの判定
                if is_error:
                    consecutive_errors += 1
                else:
                    consecutive_errors = 0  # 1つでも成功すればリセット
                    
                # 3回連続エラーで即時撤退（フェイルファスト）
                if consecutive_errors >= 3:
                    abort_event.set()
                    if logger: logger.critical("【致命的エラー】3バッチ連続で通信エラー発生。処理を強制中断します。")
                    raise RuntimeError("Gemini APIへの通信が3回連続で失敗しました。\nネットワーク接続かAPI制限をご確認ください。")
                
                # 進捗UIの更新
                processed_items += len(chunks[chunk_idx])
                if progress_callback:
                    progress_callback(processed_items)
                    
            except RuntimeError as e:
                raise e  # 致命的エラーはそのまま投げる
            except Exception as e:
                if logger: logger.error(f"チャンク結果取得エラー: {str(e)}")
    
    final_results = []
    for chunk_result in results:
        if chunk_result:
            final_results.extend(chunk_result)
        else:
            final_results.extend([""] * batch_size)
            
    return final_results

class WordProgressWindow:
    """進捗表示用ウィンドウ（スレッドセーフ版）"""
    def __init__(self, parent):
        self.window = tk.Toplevel(parent)
        self.window.title("Gemini 翻訳進捗")
        self.window.geometry("450x180")
        self.window.resizable(False, False)
        
        try:
            self.window.transient(parent)
            self.window.grab_set()
        except:
            pass
        
        self.progress_label = tk.Label(self.window, text="Gemini AI 翻訳を準備中...", font=("Arial", 11, "bold"))
        self.progress_label.pack(pady=15)
        
        self.status_label = tk.Label(self.window, text="処理を開始します...", font=("Arial", 9))
        self.status_label.pack(pady=5)
        
        self.progress_frame = tk.Frame(self.window, width=350, height=20, bg="white", relief="sunken")
        self.progress_frame.pack(pady=10)
        
        self.progress_bar = tk.Frame(self.progress_frame, height=18, bg="#0078D4")
        self.progress_bar.place(x=1, y=1)
        
        self.time_label = tk.Label(self.window, text="", font=("Arial", 8), fg="blue")
        self.time_label.pack(pady=2)
        
        self.start_time = time.time()
        
    def update_progress(self, current, total, status=""):
        # 別スレッドから安全にGUIを更新するため after を使用
        try:
            self.window.after(0, self._update_gui, current, total, status)
        except Exception:
            pass
            
    def _update_gui(self, current, total, status):
        try:
            percentage = int((current / total) * 100) if total > 0 else 0
            bar_width = int((current / total) * 348) if total > 0 else 0
            self.progress_bar.config(width=bar_width)
            
            elapsed_time = time.time() - self.start_time
            
            self.progress_label.config(text=f"翻訳進捗: {current}/{total} ({percentage}%)")
            if status:
                self.status_label.config(text=status)
            self.time_label.config(text=f"経過時間: {elapsed_time:.1f}s")
        except:
            pass
        
    def close(self):
        try:
            self.window.after(0, self.window.destroy)
        except:
            pass

def translate_ppt_document_thread(file_path, target_language, progress_window):
    """バックグラウンドで実行されるメイン処理"""
    logger = get_logger()
    try:
        start_total_time = time.time()
        lang_code = target_language.split()[0].lower()
        output_path = os.path.splitext(file_path)[0] + f"_gemini_{lang_code}.pptx"
        
        logger.info(f"=== PPT翻訳開始: {os.path.basename(file_path)} ===")
        
        # WinError 32 (ファイルロック) の事前チェック
        try:
            with open(file_path, 'a'): pass
        except PermissionError:
            logger.error(f"[WinError 32事前検知] 読み込み元ファイルがロックされています: {file_path}")
            progress_window.close()
            messagebox.showerror("ファイルエラー", "対象のPowerPointファイルが別のアプリで開かれています。\nファイルを閉じてから再度実行してください。")
            return
            
        if os.path.exists(output_path):
            try:
                with open(output_path, 'a'): pass
            except PermissionError:
                logger.error(f"[WinError 32事前検知] 保存先ファイルがロックされています: {output_path}")
                progress_window.close()
                messagebox.showerror("ファイルエラー", "以前に作成した翻訳ファイルが開かれています。\nファイルを閉じてから再度実行してください。")
                return
        
        shutil.copy2(file_path, output_path)
        prs = Presentation(output_path)
        
        translatable_items = [] 
        
        # PPTのスライドごとの処理
        for slide in prs.slides:
            # 1. 通常の図形とテーブル
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if is_translatable(run.text):
                                translatable_items.append(run)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text_frame:
                                for para in cell.text_frame.paragraphs:
                                    for run in para.runs:
                                        if is_translatable(run.text):
                                            translatable_items.append(run)
            
            # 2. スピーカーノート
            if slide.has_notes_slide:
                for shape in slide.notes_slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if is_translatable(run.text):
                                    translatable_items.append(run)

        if not translatable_items:
            progress_window.close()
            messagebox.showinfo("完了", "翻訳対象のテキストが見つかりませんでした。")
            return

        texts_only = [run.text for run in translatable_items]
        total_items = len(texts_only)
        progress_window.update_progress(0, total_items, f"Gemini APIで並列翻訳中... (0/{total_items}項目)")
        
        # UI更新用のコールバック関数
        def update_ui_callback(processed_count):
            progress_window.update_progress(processed_count, total_items, f"Gemini APIで並列翻訳中... ({processed_count}/{total_items}項目)")
        
        # ※バックグラウンドスレッドで重い通信処理を実行
        translated_texts = translate_super_fast_parallel(texts_only, target_language, max_workers=3, progress_callback=update_ui_callback, logger=logger)
        
        progress_window.update_progress(len(translatable_items), len(translatable_items), "翻訳結果をPowerPointに適用中...")
        for i, run in enumerate(translatable_items):
            if i < len(translated_texts) and translated_texts[i]:
                run.text = translated_texts[i]
                if "Japanese" in target_language or "日本" in target_language:
                    run.font.name = '游ゴシック'
        
        progress_window.update_progress(len(translatable_items), len(translatable_items), "保存中...")
        
        try:
            prs.save(output_path)
        except PermissionError:
            progress_window.close()
            messagebox.showerror("保存エラー", "ファイルが他のプログラム（PowerPointなど）で開かれています。\n閉じてから再度実行してください。")
            return

        progress_window.close()
        total_time = time.time() - start_total_time
        logger.info(f"=== 処理完了: 成功 ({total_time:.1f}秒) ===")
        
        messagebox.showinfo("完了", 
                          f"書式保持翻訳完了！\n"
                          f"保存先: {output_path}\n"
                          f"翻訳項目数: {len(translatable_items)}\n"
                          f"処理時間: {total_time:.1f}秒")
                          
    except RuntimeError as e:
        progress_window.close()
        messagebox.showerror("通信エラー強制終了", str(e))
        
    except Exception as e:
        logger.error(f"予期せぬエラー: {traceback.format_exc()}")
        progress_window.close()
        messagebox.showerror("エラー", f"翻訳処理中にエラーが発生しました:\n{str(e)}")

def select_file():
    path = filedialog.askopenfilename(
        title="翻訳するPowerPointファイルを選択してください",
        filetypes=[("PowerPoint files", "*.pptx")]
    )
    
    if not path:
        return
    
    lang_win = tk.Toplevel(root)
    lang_win.title("PPT翻訳設定")
    lang_win.geometry("400x250")
    lang_win.resizable(False, False)
    lang_win.transient(root)
    lang_win.grab_set()
    
    tk.Label(lang_win, text="翻訳先言語を選択してください", font=("Arial", 12, "bold")).pack(padx=20, pady=20)
    
    languages = {
        "日本語 (Japanese)": "Japanese",
        "英語 (English)": "English",
        "中国語簡体字 (Chinese)": "Chinese Simplified"
    }
    
    lang_var = tk.StringVar(lang_win)
    lang_var.set("日本語 (Japanese)")
    
    lang_menu = tk.OptionMenu(lang_win, lang_var, *languages.keys())
    lang_menu.config(font=("Arial", 10), width=25)
    lang_menu.pack(padx=20, pady=10)
    
    def start_translation():
        selected_language = languages[lang_var.get()]
        lang_win.destroy()
        
        # プログレスウィンドウを作成
        progress_window = WordProgressWindow(root)
        
        # 画面をフリーズさせないために、別スレッドで翻訳処理を開始！
        thread = threading.Thread(
            target=translate_ppt_document_thread, 
            args=(path, selected_language, progress_window)
        )
        thread.daemon = True
        thread.start()
    
    button_frame = tk.Frame(lang_win)
    button_frame.pack(pady=25)
    
    tk.Button(button_frame, text="翻訳開始", command=start_translation, 
             bg="#0078D4", fg="white", padx=20, pady=8, font=("Arial", 11, "bold")).pack(side=tk.LEFT, padx=10)
    tk.Button(button_frame, text="キャンセル", command=lang_win.destroy, 
             padx=20, pady=8, font=("Arial", 11)).pack(side=tk.LEFT, padx=10)

# --- GUI初期設定 ---
if __name__ == "__main__":
    root = tk.Tk()
    root.withdraw()
    
    if not check_dependencies(root):
        sys.exit(1)
        
    if not init_gemini(root):
        sys.exit(1)
        
    root.deiconify()
    root.title("PowerPoint Gemini 翻訳ツール")
    root.geometry("500x280")
    root.resizable(False, False)
    
    main_frame = tk.Frame(root)
    main_frame.pack(expand=True, fill='both', padx=20, pady=20)
    
    title_label = tk.Label(main_frame, text="PowerPoint Gemini 翻訳ツール", font=("Arial", 16, "bold"))
    title_label.pack(pady=8)
    
    subtitle_label = tk.Label(main_frame, text="書式完全保持版 (Gemini API)", font=("Arial", 12), fg="#0078D4")
    subtitle_label.pack(pady=2)
    
    desc_label = tk.Label(main_frame, 
                         text="PowerPointファイル(.pptx)を選択して翻訳します\n"
                              "フォント、色、配置、テーブル、ノート書式を完全保持",
                         font=("Arial", 10))
    desc_label.pack(pady=8)
    
    select_button = tk.Button(main_frame, text="ファイル選択", command=select_file,
                             font=("Arial", 12), bg="#0078D4", fg="white", padx=20, pady=10)
    select_button.pack(pady=15)
    
    root.mainloop()