"""ppt_translation_20260812_01.py の Gemini プロキシ移行部分の検証テスト。

Windows / tkinter / 実際の Gemini API に依存せず検証するため、
偽の `gemini_client` を **対象モジュールのロード前に** sys.modules へ注入し、
`generate_advanced()` に渡る payload を捕捉して検証する
(HANDOVER_ppt_translation_gemini_proxy.md 第8章の手法)。

python-pptx は Linux コンテナにも入るため**本物を使う**。合成PPTX(タイトル・本文・
表・スピーカーノート、run ごとにフォントサイズ/太字/色を設定)を実際に生成して
翻訳文の書き戻しまでエンドツーエンドで検証し、run 単位の書式が翻訳前後で
変わっていないこと(＝書式保持を壊していないこと)を確認する。
さらに旧版(_20260309_03)と新版に同じ翻訳文を与えて出力を比較し、
「AI呼び出し以外は何も変わっていない」ことを実物で示す。
tkinter だけスタブ化する。

実行方法:
    pip install python-pptx
    python3 tests/test_ppt_translation_20260812_01.py
"""

import importlib.util
import json
import os
import shutil
import socket
import sys
import tempfile
import types as pytypes

HERE = os.path.dirname(os.path.abspath(__file__))
TARGET = os.path.join(HERE, "..", "ppt_translation_20260812_01.py")
OLD_TARGET = os.path.join(HERE, "..", "ppt_translation_20260309_03.py")

RESULTS = []


def check(name, ok, detail=""):
    RESULTS.append((name, bool(ok), detail))
    print(f"{'PASS' if ok else 'FAIL'}: {name}" + (f"  ({detail})" if detail else ""))


# ------------------------------------------------------------
# GUI 依存のスタブ (tkinter は apt が必要でコンテナに入らない)
# ------------------------------------------------------------
def _install_env_stubs():
    class _Anything:
        """属性アクセス・呼び出しを何でも受け付けるダミー。"""
        def __init__(self, *a, **k):
            pass

        def __call__(self, *a, **k):
            return _Anything()

        def __getattr__(self, _name):
            return _Anything()

    tk = pytypes.ModuleType("tkinter")
    for attr in ("Tk", "Toplevel", "Frame", "Label", "Button", "Entry", "Checkbutton",
                 "OptionMenu", "StringVar", "BooleanVar", "LEFT", "END"):
        setattr(tk, attr, _Anything())

    filedialog = pytypes.ModuleType("tkinter.filedialog")
    filedialog.askopenfilename = lambda *a, **k: ""

    # messagebox は「呼ばれたこと」を記録して、ガードの検証に使う
    messagebox = pytypes.ModuleType("tkinter.messagebox")
    messagebox.CALLS = []

    def _record(kind):
        def _fn(title="", message="", *a, **k):
            messagebox.CALLS.append((kind, title, message))
            return True
        return _fn

    messagebox.showerror = _record("error")
    messagebox.showwarning = _record("warning")
    messagebox.showinfo = _record("info")

    tk.filedialog = filedialog
    tk.messagebox = messagebox
    sys.modules["tkinter"] = tk
    sys.modules["tkinter.filedialog"] = filedialog
    sys.modules["tkinter.messagebox"] = messagebox
    return messagebox


MESSAGEBOX = _install_env_stubs()

from pptx import Presentation  # noqa: E402  (スタブ注入後に読み込む)
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.text import PP_ALIGN  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402


# ------------------------------------------------------------
# 偽 gemini_client
# ------------------------------------------------------------
CAPTURED = {"calls": []}
RESPONSE = {"value": None}


def _numbered_reply(payload):
    """プロンプト中の [n] の数だけ番号付きの訳文を返す。"""
    text = payload["contents"][0]["parts"][0]["text"]
    n = text.count("[")
    lines = "\n".join(f"[{i}] 訳文{i}" for i in range(1, n + 1))
    return {"candidates": [{"content": {"parts": [{"text": lines}]}}],
            "usageMetadata": {"promptTokenCount": 123, "candidatesTokenCount": 45}}


def _fake_generate_advanced(payload, model=None, **kwargs):
    CAPTURED["calls"].append({"payload": payload, "model": model})
    if RESPONSE["value"] is not None:
        value = RESPONSE["value"]
        if callable(value):
            return value(payload)
        return value
    return _numbered_reply(payload)


def _load_target(module_name, inject_fake=True, target=TARGET):
    """対象ファイルをロードする。inject_fake=False なら共通モジュール未配置を再現。"""
    sys.modules.pop("gemini_client", None)
    if inject_fake:
        fake = pytypes.ModuleType("gemini_client")
        fake.generate_advanced = _fake_generate_advanced
        sys.modules["gemini_client"] = fake

    spec = importlib.util.spec_from_file_location(module_name, target)
    mod = importlib.util.module_from_spec(spec)
    sys.modules[module_name] = mod
    spec.loader.exec_module(mod)
    return mod


def _silence_sleep(mod):
    """リトライ待ち(3秒)・成功後の待ち(1.5秒)でテストが遅くならないようにする。
    対象モジュールから見える time だけを差し替え、本物の time は壊さない。"""
    stub = pytypes.ModuleType("time")
    stub.sleep = lambda _s: None
    stub.time = __import__("time").time
    mod.time = stub


# ============================================================
# 検証
# ============================================================
mod = _load_target("target_main")
_silence_sleep(mod)

check("共通モジュールを読み込めた場合 HAS_GEMINI が True", mod.HAS_GEMINI is True)
check("旧SDK(google.generativeai)への依存が無い",
      not hasattr(mod, "genai") and "google.generativeai" not in sys.modules)
check("旧SDK由来のグローバル gemini_model が残っていない",
      not hasattr(mod, "gemini_model"))
check("python-pptx への依存は残っている(HAS_PPTX)", mod.HAS_PPTX is True)


def _calls_named(path, name):
    """ソース中に「実際に呼び出している」箇所があるかをASTで調べる
    (コメント・docstringでの言及は除外する)。"""
    import ast
    tree = ast.parse(open(path, encoding="utf-8").read())
    for node in ast.walk(tree):
        if isinstance(node, ast.Call):
            fn = node.func
            if isinstance(fn, ast.Attribute) and fn.attr == name:
                return True
            if isinstance(fn, ast.Name) and fn.id == name:
                return True
    return False


check("自動モデル検出(list_models)の呼び出しが残っていない",
      not _calls_named(TARGET, "list_models"),
      "起動時のネットワークアクセスが無いこと")
check("genai.configure の呼び出しが残っていない", not _calls_named(TARGET, "configure"))
check("旧版には list_models の呼び出しがあった(移行の必要性の裏付け)",
      _calls_named(OLD_TARGET, "list_models"))

# --- init_gemini がネットワークアクセスを行わないこと -------------------
# 旧版は genai.list_models() を呼んでいたため、遮断下では必ず失敗して起動できなかった。
# ソケット生成そのものを禁止した状態で init_gemini() が成功することで、
# ネットワークアクセスが無くなったことを証明する。
os.environ["GEMINI_PROXY_URL"] = "https://example.invalid"
_real_socket = socket.socket


def _blocked_socket(*a, **k):
    raise AssertionError("init_gemini がネットワークアクセスを行った")


socket.socket = _blocked_socket
try:
    ok_init = mod.init_gemini(None)
except AssertionError as e:
    ok_init, _detail = False, str(e)
finally:
    socket.socket = _real_socket

check("init_gemini がネットワークアクセスなしで成功する", ok_init is True)
check("init_gemini がシムのクライアントを生成する",
      isinstance(mod.gemini_client, mod._CommonGeminiClient))
check("使用モデルが gemini-2.5-flash", mod.GEMINI_MODEL_NAME == "gemini-2.5-flash",
      f"model={mod.GEMINI_MODEL_NAME}")

# --- 通常の翻訳呼び出し -------------------------------------------------
CAPTURED["calls"].clear()
texts = ["Hello world", "This is a table cell"]
out, is_error = mod.translate_batch_gemini(texts, "Japanese", 0, None)

check("generate_advanced が1回呼ばれた", len(CAPTURED["calls"]) == 1,
      f"calls={len(CAPTURED['calls'])}")
check("成功時はエラーフラグが False", is_error is False)
check("response.text の解析結果が返る", out == ["訳文1", "訳文2"], f"out={out}")

call = CAPTURED["calls"][0]
payload = call["payload"]

check("model が明示的に gemini-2.5-flash で渡る", call["model"] == "gemini-2.5-flash",
      f"model={call['model']!r}")
check("payload の contents が REST 形式",
      isinstance(payload.get("contents"), list)
      and isinstance(payload["contents"][0]["parts"][0]["text"], str))
check("payload に翻訳対象テキストが [番号] 付きで含まれる",
      "[1] Hello world" in payload["contents"][0]["parts"][0]["text"]
      and "[2] This is a table cell" in payload["contents"][0]["parts"][0]["text"])
check("翻訳先言語がプロンプトに反映される",
      "Japanese" in payload["contents"][0]["parts"][0]["text"])
check("generationConfig.temperature が 0.1 で camelCase で載る",
      payload.get("generationConfig", {}).get("temperature") == 0.1,
      f"generationConfig={payload.get('generationConfig')}")
check("payload が json.dumps 可能", bool(json.dumps(payload)))
check("payload に旧SDK固有の request_options が混ざっていない",
      "request_options" not in payload and "timeout" not in payload)

# --- safetySettings (載せ忘れると応答が空になりうる) --------------------
safety = payload.get("safetySettings")
check("safetySettings が payload に載る", isinstance(safety, list), f"safetySettings={safety}")
check("safetySettings が4カテゴリ分そのまま載る",
      isinstance(safety, list) and len(safety) == 4,
      f"len={len(safety) if isinstance(safety, list) else 'N/A'}")
check("safetySettings の全カテゴリが BLOCK_NONE",
      isinstance(safety, list)
      and all(s.get("threshold") == "BLOCK_NONE" for s in safety))
check("safetySettings のカテゴリ名が REST 形式(HARM_CATEGORY_*)",
      isinstance(safety, list)
      and {s.get("category") for s in safety} == {
          "HARM_CATEGORY_HARASSMENT", "HARM_CATEGORY_HATE_SPEECH",
          "HARM_CATEGORY_SEXUALLY_EXPLICIT", "HARM_CATEGORY_DANGEROUS_CONTENT"},
      f"categories={[s.get('category') for s in safety] if isinstance(safety, list) else safety}")

# --- レスポンス互換契約 (.parts / .text / usage_metadata) ---------------
resp = mod._CommonGeminiResponse(
    {"candidates": [{"content": {"parts": [{"text": "hello"}]}}],
     "usageMetadata": {"promptTokenCount": 7, "candidatesTokenCount": 3}})
check("response.text が読める", resp.text == "hello")
check("response.parts が読める(空応答判定に必須)",
      resp.parts == [{"text": "hello"}], f"parts={resp.parts}")
check("response.usage_metadata が読める",
      resp.usage_metadata.prompt_token_count == 7
      and resp.usage_metadata.candidates_token_count == 3)

for label, raw in [("空dict", {}), ("candidates空", {"candidates": []}),
                   ("None", None), ("想定外の型", {"candidates": "broken"}),
                   ("parts空", {"candidates": [{"content": {"parts": []}}]})]:
    try:
        r = mod._CommonGeminiResponse(raw)
        ok = r.text == "" and r.parts == [] and r.usage_metadata.prompt_token_count == 0
    except Exception as e:
        ok, label = False, f"{label}: {e}"
    check(f"壊れたレスポンス({label})で例外を投げず parts が [] になる", ok)

# --- 空応答なら ValueError → 3回リトライ → 原文フォールバック -----------
RESPONSE["value"] = {"candidates": [{"content": {"parts": []}}]}
CAPTURED["calls"].clear()
out2, is_error2 = mod.translate_batch_gemini(["原文A", "原文B"], "Japanese", 1, None)
check("空応答のとき3回リトライする", len(CAPTURED["calls"]) == 3,
      f"calls={len(CAPTURED['calls'])}")
check("空応答のとき原文を返す", out2 == ["原文A", "原文B"], f"out={out2}")
check("空応答のときエラーフラグが True", is_error2 is True)
RESPONSE["value"] = None

# --- 通信失敗時も3回リトライして原文フォールバック ----------------------
def _raising(payload):
    raise RuntimeError("proxy down")


RESPONSE["value"] = _raising
CAPTURED["calls"].clear()
out3, is_error3 = mod.translate_batch_gemini(["原文A"], "Japanese", 2, None)
check("通信失敗時も3回リトライする", len(CAPTURED["calls"]) == 3,
      f"calls={len(CAPTURED['calls'])}")
check("通信失敗時は例外を投げず原文を返す", out3 == ["原文A"], f"out={out3}")
check("通信失敗時はエラーフラグが True", is_error3 is True)
RESPONSE["value"] = None

# --- 番号が一部欠けた応答は該当項目だけ原文 -----------------------------
RESPONSE["value"] = {"candidates": [{"content": {"parts": [{"text": "[1] OK-1"}]}}]}
out4, _ = mod.translate_batch_gemini(["原文A", "原文B"], "Japanese", 3, None)
check("欠番の項目だけ原文にフォールバックする", out4 == ["OK-1", "原文B"], f"out={out4}")
RESPONSE["value"] = None

# --- 3バッチ連続エラーでフェイルファストする(移行後も維持されているか) --
RESPONSE["value"] = _raising
try:
    mod.translate_super_fast_parallel(["a" * 5] * 40, "Japanese", max_workers=1, logger=None)
    check("3バッチ連続エラーでフェイルファストする", False, "RuntimeErrorが出なかった")
except RuntimeError as e:
    check("3バッチ連続エラーでフェイルファストする", True)
    check("フェイルファストのメッセージが従来どおり", "3回連続で失敗" in str(e), str(e))
RESPONSE["value"] = None

# --- 並列翻訳が正常系で全件返すこと -------------------------------------
res = mod.translate_super_fast_parallel(["text %d" % i for i in range(25)],
                                        "Japanese", max_workers=3, logger=None)
check("並列翻訳が入力と同じ件数を返す", len(res) == 25, f"len={len(res)}")

# --- 進捗コールバックが従来どおり呼ばれる -------------------------------
_progress = []
mod.translate_super_fast_parallel(["text %d" % i for i in range(25)], "Japanese",
                                  max_workers=3,
                                  progress_callback=_progress.append, logger=None)
check("進捗コールバックがバッチごとに呼ばれる", len(_progress) == 3, f"calls={_progress}")
check("進捗コールバックの最終値が全項目数になる",
      _progress and max(_progress) == 25, f"max={max(_progress) if _progress else None}")

# --- is_translatable (移行と無関係だが回帰確認) -------------------------
check("is_translatable: 通常の文は翻訳対象", mod.is_translatable("Hello world") is True)
check("is_translatable: 空文字は対象外", mod.is_translatable("") is False)
check("is_translatable: 2文字以下は対象外", mod.is_translatable("ab") is False)
check("is_translatable: 数字だけは対象外", mod.is_translatable("12.34") is False)
check("is_translatable: 箇条書き記号は対象外", mod.is_translatable("•") is False)

# --- 認証情報の判定 (プロキシURLのみ = 通ることが最重要) ----------------
_env_backup = {k: os.environ.get(k) for k in ("GEMINI_API_KEY", "GEMINI_PROXY_URL")}
for api, proxy, expected, label in [
        ("k", None, True, "APIキーのみ"),
        (None, "https://x", True, "プロキシURLのみ(プロキシ専用構成)"),
        ("k", "https://x", True, "両方"),
        (None, None, False, "どちらも未設定")]:
    for name, val in (("GEMINI_API_KEY", api), ("GEMINI_PROXY_URL", proxy)):
        os.environ.pop(name, None)
        if val:
            os.environ[name] = val
    check(f"gemini_credentials_available: {label} -> {expected}",
          mod.gemini_credentials_available() is expected)

    MESSAGEBOX.CALLS.clear()
    check(f"init_gemini: {label} -> {expected}", mod.init_gemini(None) is expected)
    if not expected:
        msg = MESSAGEBOX.CALLS[-1][2] if MESSAGEBOX.CALLS else ""
        check("未設定時のエラーが両方の環境変数を案内する",
              "GEMINI_API_KEY" in msg and "GEMINI_PROXY_URL" in msg, msg)

for k, v in _env_backup.items():
    os.environ.pop(k, None)
    if v:
        os.environ[k] = v

# --- 共通モジュール未配置時 ---------------------------------------------
mod2 = _load_target("target_no_common", inject_fake=False)
check("共通モジュール未配置なら HAS_GEMINI が False", mod2.HAS_GEMINI is False)

try:
    mod2._CommonGeminiClient().models.generate_content(model="m", contents="c")
    check("共通モジュール未配置でAI呼び出しすると RuntimeError", False, "例外が出なかった")
except RuntimeError as e:
    text = str(e)
    check("共通モジュール未配置でAI呼び出しすると RuntimeError", True)
    check("エラーに探索したパスが含まれる", "common" in text, text.splitlines()[1])
    check("エラーに GEMINI_COMMON_DIR の案内が含まれる", "GEMINI_COMMON_DIR" in text)
except Exception as e:
    check("共通モジュール未配置でAI呼び出しすると RuntimeError", False, repr(e))

MESSAGEBOX.CALLS.clear()
check("共通モジュール未配置なら check_dependencies が False",
      mod2.check_dependencies(None) is False)
check("check_dependencies のエラーが原因を説明している",
      any("gemini_client.py" in c[2] for c in MESSAGEBOX.CALLS),
      f"calls={MESSAGEBOX.CALLS}")
MESSAGEBOX.CALLS.clear()
check("共通モジュール未配置なら init_gemini も False", mod2.init_gemini(None) is False)
check("依存関係が揃っていれば check_dependencies が True",
      mod.check_dependencies(None) is True)

# --- 共通モジュールの探索 (会社PCは「2つ上」が common) ------------------
_tmp = tempfile.mkdtemp()
try:
    _common = os.path.join(_tmp, "PythonScripts", "common")
    _tool = os.path.join(_tmp, "PythonScripts", "Powerpoint", "ppt_translator")
    os.makedirs(_common)
    os.makedirs(_tool)
    with open(os.path.join(_common, "gemini_client.py"), "w", encoding="utf-8") as f:
        f.write("MARKER = 'two-up'\n\n\ndef generate_advanced(payload, model=None):\n"
                "    return {'candidates': [{'content': {'parts': [{'text': '[1] ok'}]}}]}\n")
    _copied = os.path.join(_tool, os.path.basename(TARGET))
    shutil.copy(TARGET, _copied)

    _saved_path = list(sys.path)
    mod3 = _load_target("target_two_up", inject_fake=False, target=_copied)
    check("「2つ上が common」の配置で共通モジュールを解決できる",
          mod3.HAS_GEMINI is True and mod3._COMMON_DIR == os.path.abspath(_common),
          f"COMMON_DIR={getattr(mod3, '_COMMON_DIR', None)}")
    check("探索候補に1つ上・2つ上・3つ上が含まれる",
          len(mod3._COMMON_DIR_CANDIDATES) == 3, f"{mod3._COMMON_DIR_CANDIDATES}")

    # 移行期間中に元の場所(PythonScripts\excel\)へ置いても動くこと(＝「1つ上」)
    _excel = os.path.join(_tmp, "PythonScripts", "excel")
    os.makedirs(_excel)
    _copied2 = os.path.join(_excel, os.path.basename(TARGET))
    shutil.copy(TARGET, _copied2)
    mod3b = _load_target("target_one_up", inject_fake=False, target=_copied2)
    check("「1つ上が common」の配置(移行元 excel フォルダ)でも解決できる",
          mod3b.HAS_GEMINI is True and mod3b._COMMON_DIR == os.path.abspath(_common),
          f"COMMON_DIR={getattr(mod3b, '_COMMON_DIR', None)}")

    # GEMINI_COMMON_DIR による明示指定
    os.environ["GEMINI_COMMON_DIR"] = _common
    mod4 = _load_target("target_env_dir", inject_fake=False)
    check("GEMINI_COMMON_DIR で共通モジュールの場所を明示指定できる",
          mod4.HAS_GEMINI is True
          and mod4._COMMON_DIR_CANDIDATES == [_common],
          f"candidates={mod4._COMMON_DIR_CANDIDATES}")
    os.environ.pop("GEMINI_COMMON_DIR", None)
    sys.path[:] = _saved_path
finally:
    shutil.rmtree(_tmp, ignore_errors=True)
    sys.modules.pop("gemini_client", None)

# --- 未使用オプション(他ツールとの契約統一) -----------------------------
CAPTURED["calls"].clear()
sys.modules["gemini_client"] = pytypes.ModuleType("gemini_client")
cfg = mod._GeminiGenerateConfig(temperature=0.5, response_mime_type="application/json",
                                response_schema={"type": "object"},
                                system_instruction="be brief")
mod._CommonGeminiClient().models.generate_content(model="m", contents="c", config=cfg)
_p = CAPTURED["calls"][-1]["payload"]
check("responseMimeType / responseSchema が camelCase で載る",
      _p.get("generationConfig", {}).get("responseMimeType") == "application/json"
      and _p.get("generationConfig", {}).get("responseSchema") == {"type": "object"},
      f"generationConfig={_p.get('generationConfig')}")
check("systemInstruction が REST 形式で載る",
      _p.get("systemInstruction") == {"parts": [{"text": "be brief"}]},
      f"systemInstruction={_p.get('systemInstruction')}")

CAPTURED["calls"].clear()
mod._CommonGeminiClient().models.generate_content(model="m", contents="c", config=None)
check("config が None でも payload を組み立てられる",
      "generationConfig" not in CAPTURED["calls"][-1]["payload"]
      and "safetySettings" not in CAPTURED["calls"][-1]["payload"])


class _FakePydanticSchema:
    def model_dump(self, **kwargs):
        return {"type": "OBJECT", "propertyOrdering": ["a"]}


check("pydantic モデルの response_schema も dict 化できる",
      mod._schema_to_jsonable(_FakePydanticSchema()) == {"type": "OBJECT",
                                                         "propertyOrdering": ["a"]})


# ============================================================
# PowerPoint処理のエンドツーエンド検証 (python-pptx は本物を使う)
#   移行で PPT の書式保持処理に一切影響が出ていないことを、合成PPTXで確認する。
#   run 単位の書式(サイズ/太字/色/段落配置)が翻訳前後で一致すれば、
#   「書式完全保持」を壊していない証拠になる。
# ============================================================
class _FakeProgressWindow:
    """WordProgressWindow 互換の記録用ダミー(tkinter を使わない)。"""
    def __init__(self):
        self.updates = []
        self.closed = 0

    def update_progress(self, current, total, status=""):
        self.updates.append((current, total, status))

    def close(self):
        self.closed += 1


def _build_sample_pptx(path):
    """タイトル ＋ 本文(複数run・書式バラバラ) ＋ 表 ＋ スピーカーノート を含むPPTXを作る。"""
    prs = Presentation()
    blank = prs.slide_layouts[6]

    # --- スライド1: タイトルと本文 ---
    slide = prs.slides.add_slide(blank)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
    tp = title_box.text_frame.paragraphs[0]
    tp.alignment = PP_ALIGN.CENTER
    r = tp.add_run()
    r.text = "Employee Engagement Survey Results"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0xC0, 0x30, 0x30)

    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9), Inches(2))
    p1 = body.text_frame.paragraphs[0]
    p1.alignment = PP_ALIGN.LEFT
    r1 = p1.add_run()
    r1.text = "This paragraph describes the overall findings."
    r1.font.size = Pt(14)
    r1.font.bold = False
    r1.font.color.rgb = RGBColor(0x20, 0x20, 0x20)
    # 同じ段落の2つ目の run (書式が違う = run 単位保持の検証に使う)
    r2 = p1.add_run()
    r2.text = " Response rate reached a record high."
    r2.font.size = Pt(11)
    r2.font.bold = True
    r2.font.color.rgb = RGBColor(0x00, 0x66, 0xCC)

    p2 = body.text_frame.add_paragraph()
    p2.alignment = PP_ALIGN.RIGHT
    r3 = p2.add_run()
    r3.text = "Prepared by the strategy office."
    r3.font.size = Pt(10)
    r3.font.color.rgb = RGBColor(0x80, 0x80, 0x80)

    # 翻訳対象外になるはずの run (2文字以下 / 記号 / 数字)
    small = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(3), Inches(0.8))
    sp = small.text_frame.paragraphs[0]
    for txt, size in (("ab", 12), ("•", 12), ("12.34", 12)):
        rr = sp.add_run()
        rr.text = txt
        rr.font.size = Pt(size)

    # --- スライド2: 表 ---
    slide2 = prs.slides.add_slide(blank)
    table_shape = slide2.shapes.add_table(3, 2, Inches(0.5), Inches(0.5),
                                          Inches(8), Inches(2))
    table = table_shape.table
    cells = [("Category", "Score"), ("Wellbeing", "Sixty six percent"),
             ("Leadership", "Seventy two percent")]
    for i, (left, right) in enumerate(cells):
        for j, txt in enumerate((left, right)):
            cell = table.cell(i, j)
            para = cell.text_frame.paragraphs[0]
            run = para.add_run()
            run.text = txt
            run.font.size = Pt(12 + i)
            run.font.bold = (i == 0)
            run.font.color.rgb = RGBColor(0x11, 0x22, 0x33)

    # --- スピーカーノート ---
    notes = slide.notes_slide.notes_text_frame
    notes.text = "Remember to mention the response rate improvement."

    prs.save(path)


def _run_signature(path):
    """run 単位のテキストと書式を、比較できる形で書き出す。"""
    prs = Presentation(path)
    sig = []

    def _walk(text_frame, key):
        for pi, para in enumerate(text_frame.paragraphs):
            for ri, run in enumerate(para.runs):
                try:
                    color = str(run.font.color.rgb)
                except Exception:
                    color = None
                sig.append({
                    "key": key + (pi, ri),
                    "text": run.text,
                    "size": run.font.size.pt if run.font.size is not None else None,
                    "bold": run.font.bold,
                    "name": run.font.name,
                    "color": color,
                    "align": str(para.alignment),
                })

    for si, slide in enumerate(prs.slides):
        for shi, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                _walk(shape.text_frame, (si, shi))
            if shape.has_table:
                for ri_, row in enumerate(shape.table.rows):
                    for ci, cell in enumerate(row.cells):
                        _walk(cell.text_frame, (si, shi, ri_, ci))
        if slide.has_notes_slide:
            for shi, shape in enumerate(slide.notes_slide.shapes):
                if shape.has_text_frame:
                    _walk(shape.text_frame, (si, "notes", shi))
    return sig


def _fmt_only(sig):
    """書式だけを抜き出す(テキストとフォント名は翻訳で変わるので除外)。"""
    return [{k: v for k, v in item.items() if k not in ("text", "name")} for item in sig]


_ppt_dir = tempfile.mkdtemp()
_cwd = os.getcwd()
try:
    # translation_debug.log がリポジトリに落ちないよう、一時フォルダで実行する
    os.chdir(_ppt_dir)

    src_pptx = os.path.join(_ppt_dir, "sample.pptx")
    _build_sample_pptx(src_pptx)
    before_sig = _run_signature(src_pptx)

    CAPTURED["calls"].clear()
    pw = _FakeProgressWindow()
    MESSAGEBOX.CALLS.clear()
    mod.translate_ppt_document_thread(src_pptx, "Japanese", pw)

    out_pptx = os.path.join(_ppt_dir, "sample_gemini_japanese.pptx")
    check("出力ファイル名が 元ファイル名_gemini_japanese.pptx になる",
          os.path.isfile(out_pptx), f"exists={os.path.isfile(out_pptx)}")
    check("エラーダイアログが出ていない",
          not any(c[0] == "error" for c in MESSAGEBOX.CALLS), f"{MESSAGEBOX.CALLS}")
    check("完了ダイアログが出る", any(c[0] == "info" for c in MESSAGEBOX.CALLS))
    check("進捗ウィンドウが閉じられる", pw.closed >= 1, f"closed={pw.closed}")
    check("進捗ウィンドウが更新される", len(pw.updates) > 0, f"updates={len(pw.updates)}")
    check("翻訳がシム経由(共通モジュール)で行われた", len(CAPTURED["calls"]) >= 1,
          f"calls={len(CAPTURED['calls'])}")
    check("エンドツーエンドでも model が明示的に渡る",
          all(c["model"] == "gemini-2.5-flash" for c in CAPTURED["calls"]))
    check("エンドツーエンドでも safetySettings が載る",
          all(len(c["payload"].get("safetySettings", [])) == 4 for c in CAPTURED["calls"]))

    after_sig = _run_signature(out_pptx)

    check("run の構成(数と位置)が翻訳前後で変わらない",
          [s["key"] for s in before_sig] == [s["key"] for s in after_sig],
          f"before={len(before_sig)} / after={len(after_sig)}")
    check("run 単位の書式(サイズ/太字/色/段落配置)が翻訳前後で完全一致する",
          _fmt_only(before_sig) == _fmt_only(after_sig),
          "差分あり" if _fmt_only(before_sig) != _fmt_only(after_sig) else "")

    _translated = [s for s in after_sig if s["text"].startswith("訳文")]
    check("翻訳文が run に書き戻されている", len(_translated) > 0,
          f"translated={len(_translated)}")
    check("日本語指定時に run.font.name が 游ゴシック になる",
          all(s["name"] == "游ゴシック" for s in _translated),
          f"names={sorted({s['name'] for s in _translated})}")

    # 表・ノートも翻訳対象として拾えていること
    check("表のセルが翻訳対象として拾われている",
          any(len(s["key"]) == 6 and s["text"].startswith("訳文") for s in after_sig),
          f"table_runs={[s['key'] for s in after_sig if len(s['key']) == 6]}")
    check("スピーカーノートが翻訳対象として拾われている",
          any("notes" in s["key"] and s["text"].startswith("訳文") for s in after_sig))

    # 翻訳対象外のものは原文のまま
    _untouched = {s["key"]: s["text"] for s in after_sig}
    check("2文字以下・記号・数字だけの run は原文のまま",
          all(t in _untouched.values() for t in ("ab", "•", "12.34")),
          f"texts={[s['text'] for s in after_sig]}")

    # --- 旧版(_20260309_03)と新版でPPTX出力が一致することの直接確認 --------
    # 移行が触ったのは Gemini 呼び出し経路だけなので、同じ翻訳文を与えれば
    # 旧版と新版のPPTX出力(テキスト・書式)は完全に一致するはずである。
    # ※PPTX(zip)はタイムスタンプが毎回変わるためバイト比較はできない。
    #   run 単位のテキストと書式で比較する。
    if os.path.isfile(OLD_TARGET):
        class _OldResponse:
            def __init__(self, payload_text):
                n = payload_text.count("[")
                self.text = "\n".join(f"[{i}] 訳文{i}" for i in range(1, n + 1))
                self.parts = [{"text": self.text}]

        class _OldModel:
            def generate_content(self, prompt, **kwargs):
                return _OldResponse(prompt)

        fake_genai = pytypes.ModuleType("google.generativeai")
        fake_genai.configure = lambda **k: None
        fake_genai.list_models = lambda: []
        fake_genai.GenerativeModel = lambda *a, **k: _OldModel()
        fake_genai.types = pytypes.SimpleNamespace(
            GenerationConfig=lambda **k: pytypes.SimpleNamespace(**k))
        google_pkg = pytypes.ModuleType("google")
        google_pkg.generativeai = fake_genai
        sys.modules["google"] = google_pkg
        sys.modules["google.generativeai"] = fake_genai

        spec = importlib.util.spec_from_file_location("target_old", OLD_TARGET)
        old_mod = importlib.util.module_from_spec(spec)
        sys.modules["target_old"] = old_mod
        spec.loader.exec_module(old_mod)
        _silence_sleep(old_mod)
        old_mod.gemini_model = _OldModel()

        old_src = os.path.join(_ppt_dir, "sample_old.pptx")
        shutil.copy(src_pptx, old_src)
        old_mod.translate_ppt_document_thread(old_src, "Japanese", _FakeProgressWindow())
        old_out = os.path.join(_ppt_dir, "sample_old_gemini_japanese.pptx")

        check("旧版でも同じ名前で出力される", os.path.isfile(old_out))
        old_sig = _run_signature(old_out)
        check("旧版と新版で出力PPTXの run 構成が一致する",
              [s["key"] for s in old_sig] == [s["key"] for s in after_sig],
              f"old={len(old_sig)} / new={len(after_sig)}")
        check("旧版と新版で出力PPTXのテキストが完全一致する",
              [s["text"] for s in old_sig] == [s["text"] for s in after_sig])
        check("旧版と新版で出力PPTXの書式が完全一致する",
              _fmt_only(old_sig) == _fmt_only(after_sig))
        check("旧版と新版で出力PPTXのフォント名が完全一致する",
              [s["name"] for s in old_sig] == [s["name"] for s in after_sig])
    else:
        check("旧版(_20260309_03)が見つかり比較できる", False, f"not found: {OLD_TARGET}")

    # --- 翻訳対象が無いPPTXでも落ちないこと -------------------------------
    empty_pptx = os.path.join(_ppt_dir, "empty.pptx")
    Presentation().save(empty_pptx)
    MESSAGEBOX.CALLS.clear()
    pw2 = _FakeProgressWindow()
    mod.translate_ppt_document_thread(empty_pptx, "Japanese", pw2)
    check("翻訳対象が無いPPTXでもエラーにならず案内が出る",
          any(c[0] == "info" for c in MESSAGEBOX.CALLS)
          and not any(c[0] == "error" for c in MESSAGEBOX.CALLS),
          f"{MESSAGEBOX.CALLS}")
finally:
    os.chdir(_cwd)
    shutil.rmtree(_ppt_dir, ignore_errors=True)
    sys.modules.pop("google", None)
    sys.modules.pop("google.generativeai", None)


# ============================================================
print("\n" + "=" * 60)
failed = [r for r in RESULTS if not r[1]]
print(f"合計 {len(RESULTS)} 項目 / 合格 {len(RESULTS) - len(failed)} / 失敗 {len(failed)}")
if failed:
    for name, _ok, detail in failed:
        print(f"  FAILED: {name}  {detail}")
sys.exit(1 if failed else 0)
